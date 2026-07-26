import io

import fitz
from PIL import Image
from pptx import Presentation

from app import app


def make_image(size=(800, 600), color=(230, 100, 30), image_format="PNG"):
    stream = io.BytesIO()
    Image.new("RGB", size, color).save(stream, format=image_format)
    stream.seek(0)
    return stream


def make_pdf(page_count=2):
    document = fitz.open()
    for page_number in range(page_count):
        page = document.new_page(width=842, height=595)
        page.insert_text((72, 72), f"Testseite {page_number + 1}")
    stream = io.BytesIO(document.tobytes())
    document.close()
    stream.seek(0)
    return stream


def test_health_endpoint():
    with app.test_client() as client:
        response = client.get("/healthz")

    assert response.status_code == 200
    assert response.get_json() == {"status": "ok"}


def test_image_and_pdf_create_one_slide_per_page():
    with app.test_client() as client:
        response = client.post(
            "/upload",
            data={
                "files": [
                    (make_image(), "bild.png"),
                    (make_pdf(page_count=2), "folien.pdf"),
                ]
            },
            content_type="multipart/form-data",
        )
        presentation_bytes = response.get_data()
        response.close()

    assert response.status_code == 200
    assert response.mimetype == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    presentation = Presentation(io.BytesIO(presentation_bytes))
    assert len(presentation.slides) == 3
    assert all(len(slide.shapes) == 1 for slide in presentation.slides)


def test_portrait_image_is_centered_without_distortion():
    with app.test_client() as client:
        response = client.post(
            "/upload",
            data={"files": (make_image(size=(600, 900)), "portrait.jpg")},
            content_type="multipart/form-data",
        )
        presentation_bytes = response.get_data()
        response.close()

    presentation = Presentation(io.BytesIO(presentation_bytes))
    picture = presentation.slides[0].shapes[0]
    assert picture.height == presentation.slide_height
    assert picture.width < presentation.slide_width
    assert picture.left == (presentation.slide_width - picture.width) // 2


def test_broken_image_returns_clear_error_instead_of_partial_pptx():
    with app.test_client() as client:
        response = client.post(
            "/upload",
            data={"files": (io.BytesIO(b"not an image"), "defekt.png")},
            content_type="multipart/form-data",
        )

    assert response.status_code == 400
    assert "defekt.png" in response.get_json()["error"]


def test_unsupported_extension_is_rejected():
    with app.test_client() as client:
        response = client.post(
            "/upload",
            data={"files": (io.BytesIO(b"text"), "notizen.txt")},
            content_type="multipart/form-data",
        )

    assert response.status_code == 400
    assert "nicht unterstützten Dateityp" in response.get_json()["error"]


def test_pdf_page_limit_is_enforced():
    previous_limit = app.config["MAX_PDF_PAGES"]
    app.config["MAX_PDF_PAGES"] = 1
    try:
        with app.test_client() as client:
            response = client.post(
                "/upload",
                data={"files": (make_pdf(page_count=2), "zu-viele-seiten.pdf")},
                content_type="multipart/form-data",
            )
    finally:
        app.config["MAX_PDF_PAGES"] = previous_limit

    assert response.status_code == 400
    assert "höchstens 1 Seite" in response.get_json()["error"]


def test_total_slide_limit_is_enforced_across_files():
    previous_limit = app.config["MAX_SLIDES"]
    app.config["MAX_SLIDES"] = 1
    try:
        with app.test_client() as client:
            response = client.post(
                "/upload",
                data={
                    "files": [
                        (make_image(), "eins.png"),
                        (make_image(), "zwei.png"),
                    ]
                },
                content_type="multipart/form-data",
            )
    finally:
        app.config["MAX_SLIDES"] = previous_limit

    assert response.status_code == 400
    assert "höchstens 1 Folie" in response.get_json()["error"]
