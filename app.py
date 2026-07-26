import io
import logging
import os
import tempfile
import warnings
from pathlib import Path

import fitz
from flask import Flask, jsonify, render_template, request, send_file
from PIL import Image, ImageOps, UnidentifiedImageError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu


SLIDE_WIDTH_EMU = 9_144_000
SLIDE_HEIGHT_EMU = 5_143_500
MAX_IMG_WIDTH_PX = 1_920
MAX_IMG_HEIGHT_PX = 1_080
ALLOWED_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".pdf"}

Image.MAX_IMAGE_PIXELS = 25_000_000

app = Flask(__name__)
app.config.from_mapping(
    MAX_CONTENT_LENGTH=int(os.getenv("MAX_UPLOAD_MB", "100")) * 1024 * 1024,
    MAX_FILES=int(os.getenv("MAX_FILES", "250")),
    MAX_PDF_PAGES=int(os.getenv("MAX_PDF_PAGES", "250")),
    MAX_SLIDES=int(os.getenv("MAX_SLIDES", "250")),
)

logging.basicConfig(level=os.getenv("LOG_LEVEL", "INFO"))
logger = logging.getLogger(__name__)


class ConversionError(ValueError):
    """Fehler, der dem Benutzer als ungültige Eingabe gemeldet werden kann."""


def _normalized_image_stream(image_source) -> io.BytesIO:
    """Liest ein Bild, korrigiert EXIF und liefert ein begrenztes PNG."""
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("error", Image.DecompressionBombWarning)
            with Image.open(image_source) as source:
                source.seek(0)
                source_format = source.format
                image = ImageOps.exif_transpose(source)
                image.load()
                image.thumbnail(
                    (MAX_IMG_WIDTH_PX, MAX_IMG_HEIGHT_PX),
                    Image.Resampling.LANCZOS,
                )

                has_alpha = image.mode in {"RGBA", "LA"} or "transparency" in image.info
                target_mode = "RGBA" if has_alpha else "RGB"
                if image.mode != target_mode:
                    image = image.convert(target_mode)

                output = io.BytesIO()
                if source_format == "JPEG" and not has_alpha:
                    image.save(output, format="JPEG", quality=92, optimize=True)
                else:
                    image.save(output, format="PNG")
                output.seek(0)
                return output
    except (
        UnidentifiedImageError,
        Image.DecompressionBombWarning,
        Image.DecompressionBombError,
        OSError,
        ValueError,
    ) as exc:
        raise ConversionError("Die Bilddatei ist beschädigt oder nicht unterstützt.") from exc


def _add_image_centered(slide, image_stream: io.BytesIO) -> None:
    """Fügt ein Bild proportional und zentriert auf einer Folie ein."""
    image_stream.seek(0)
    with Image.open(image_stream) as image:
        width_px, height_px = image.size

    if width_px <= 0 or height_px <= 0:
        raise ConversionError("Ein Bild besitzt ungültige Abmessungen.")

    scale = min(SLIDE_WIDTH_EMU / width_px, SLIDE_HEIGHT_EMU / height_px)
    width_emu = int(width_px * scale)
    height_emu = int(height_px * scale)
    left = int((SLIDE_WIDTH_EMU - width_emu) / 2)
    top = int((SLIDE_HEIGHT_EMU - height_emu) / 2)

    image_stream.seek(0)
    slide.shapes.add_picture(
        image_stream,
        left,
        top,
        width=width_emu,
        height=height_emu,
    )


def _add_slide(prs: Presentation, image_stream: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    _add_image_centered(slide, image_stream)


def _render_pdf_pages(pdf_path: Path):
    try:
        document = fitz.open(pdf_path)
    except (fitz.FileDataError, RuntimeError) as exc:
        raise ConversionError("Die PDF-Datei ist beschädigt oder nicht unterstützt.") from exc

    try:
        if document.needs_pass:
            raise ConversionError("Passwortgeschützte PDFs werden nicht unterstützt.")
        if document.page_count == 0:
            raise ConversionError("Die PDF-Datei enthält keine Seiten.")
        if document.page_count > app.config["MAX_PDF_PAGES"]:
            page_word = "Seite" if app.config["MAX_PDF_PAGES"] == 1 else "Seiten"
            raise ConversionError(
                f"Eine PDF darf höchstens {app.config['MAX_PDF_PAGES']} "
                f"{page_word} enthalten."
            )

        for page in document:
            page_width = max(page.rect.width, 1)
            page_height = max(page.rect.height, 1)
            zoom = min(
                MAX_IMG_WIDTH_PX / page_width,
                MAX_IMG_HEIGHT_PX / page_height,
            )
            pixmap = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
            stream = io.BytesIO(pixmap.tobytes("png"))
            stream.seek(0)
            yield stream
    except (fitz.FileDataError, RuntimeError) as exc:
        raise ConversionError("Eine PDF-Seite konnte nicht gelesen werden.") from exc
    finally:
        document.close()


def _convert_uploads(files, workdir: Path, output_path: Path) -> int:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)
    slide_count = 0

    for index, uploaded_file in enumerate(files):
        display_name = uploaded_file.filename or f"Datei {index + 1}"
        extension = Path(display_name).suffix.lower()
        if extension not in ALLOWED_EXTENSIONS:
            raise ConversionError(f"„{display_name}“ hat einen nicht unterstützten Dateityp.")

        upload_path = workdir / f"upload-{index}{extension}"
        uploaded_file.save(upload_path)

        try:
            if extension == ".pdf":
                for page_stream in _render_pdf_pages(upload_path):
                    try:
                        if slide_count >= app.config["MAX_SLIDES"]:
                            slide_word = (
                                "Folie" if app.config["MAX_SLIDES"] == 1 else "Folien"
                            )
                            raise ConversionError(
                                "Die Präsentation darf insgesamt höchstens "
                                f"{app.config['MAX_SLIDES']} {slide_word} enthalten."
                            )
                        _add_slide(prs, page_stream)
                        slide_count += 1
                    finally:
                        page_stream.close()
            else:
                image_stream = _normalized_image_stream(upload_path)
                try:
                    if slide_count >= app.config["MAX_SLIDES"]:
                        slide_word = (
                            "Folie" if app.config["MAX_SLIDES"] == 1 else "Folien"
                        )
                        raise ConversionError(
                            "Die Präsentation darf insgesamt höchstens "
                            f"{app.config['MAX_SLIDES']} {slide_word} enthalten."
                        )
                    _add_slide(prs, image_stream)
                    slide_count += 1
                finally:
                    image_stream.close()
        except ConversionError as exc:
            raise ConversionError(f"Fehler in „{display_name}“: {exc}") from exc

    if slide_count == 0:
        raise ConversionError("Es wurden keine gültigen Seiten oder Bilder gefunden.")

    prs.save(output_path)
    return slide_count


@app.errorhandler(413)
def upload_too_large(_error):
    max_mb = app.config["MAX_CONTENT_LENGTH"] // (1024 * 1024)
    return jsonify(error=f"Der Upload darf insgesamt höchstens {max_mb} MB groß sein."), 413


@app.get("/")
def index():
    return render_template("index.html")


@app.get("/healthz")
def health():
    return jsonify(status="ok")


@app.post("/upload")
def upload_files():
    files = [file for file in request.files.getlist("files") if file.filename]
    if not files:
        return jsonify(error="Keine Dateien ausgewählt."), 400
    if len(files) > app.config["MAX_FILES"]:
        return jsonify(
            error=f"Es dürfen höchstens {app.config['MAX_FILES']} Dateien hochgeladen werden."
        ), 400

    output_file = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    output_path = Path(output_file.name)
    output_file.close()

    try:
        with tempfile.TemporaryDirectory(prefix="2pptx-") as temp_dir:
            slide_count = _convert_uploads(files, Path(temp_dir), output_path)

        logger.info(
            "Präsentation mit %d Folien aus %d Dateien erstellt.",
            slide_count,
            len(files),
        )
        response = send_file(
            output_path,
            as_attachment=True,
            download_name="presentation.pptx",
            mimetype=(
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            ),
        )
        response.call_on_close(lambda: output_path.unlink(missing_ok=True))
        return response
    except ConversionError as exc:
        output_path.unlink(missing_ok=True)
        logger.warning("Konvertierung abgelehnt: %s", exc)
        return jsonify(error=str(exc)), 400
    except Exception:
        output_path.unlink(missing_ok=True)
        logger.exception("Unerwarteter Fehler bei der Konvertierung")
        return jsonify(error="Die Präsentation konnte nicht erstellt werden."), 500


if __name__ == "__main__":
    app.run(debug=False, host="0.0.0.0", port=5000)
